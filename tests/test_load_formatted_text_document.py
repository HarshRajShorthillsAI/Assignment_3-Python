import pytest
from docx import Document
from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
from loaders.ppt_loader import PPTLoader

from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.shapes import MSO_SHAPE

import os

from extractors.data_extractor import DataExtractor


from fpdf import FPDF
from fpdf.enums import XPos, YPos
from PIL import Image
import os


class TestFormattedTextLoading:
    FORMATTED_CONTENT = {
        "bold_text": "This text is bold.",
        "italic_text": "This text is italic.",
        "underlined_text": "This text is underlined.",
        "colored_text": "This text is in red color.",
        "large_text": "This text is large.",
        "centered_text": "This text is centered.",
        "heading1_text": "This is Heading 1",
        "heading2_text": "This is Heading 2",
        "hyperlink_text": "OpenAI website",
        "list_items": ["First item", "Second item", "Third item"],
        "table_data": [["Header 1", "Header 2"], ["Row 1, Cell 1", "Row 1, Cell 2"], ["Row 2, Cell 1", "Row 2, Cell 2"]],
        "image_path": "tests/test_files/img.png"
    }

    class PDF(FPDF):
        def header(self):
            self.set_font("Helvetica", 'B', 12)
            self.cell(0, 10, 'Title', new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')

        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", 'I', 8)
            self.cell(0, 10, f'Page {self.page_no()}', new_x=XPos.RIGHT, new_y=YPos.NEXT, align='C')


    @staticmethod
    def create_formatted_docx_file(folder_path, filename="formatted_test_file.docx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        # Create a DOCX document
        doc = Document()

        # Bold text
        bold_paragraph = doc.add_paragraph()
        bold_run = bold_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["bold_text"])
        bold_run.bold = True

        # Italic text
        italic_paragraph = doc.add_paragraph()
        italic_run = italic_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["italic_text"])
        italic_run.italic = True

        # Underlined text
        underline_paragraph = doc.add_paragraph()
        underline_run = underline_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["underlined_text"])
        underline_run.underline = True

        # Colored text
        color_paragraph = doc.add_paragraph()
        color_run = color_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["colored_text"])
        color_run.font.color.rgb = RGBColor(255, 0, 0)  # Red color

        # Large text
        large_text_paragraph = doc.add_paragraph()
        large_text_run = large_text_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["large_text"])
        large_text_run.font.size = Pt(18)

        # Centered text
        centered_paragraph = doc.add_paragraph()
        centered_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        centered_run = centered_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["centered_text"])

        # Headings
        doc.add_heading(TestFormattedTextLoading.FORMATTED_CONTENT["heading1_text"], level=1)
        doc.add_heading(TestFormattedTextLoading.FORMATTED_CONTENT["heading2_text"], level=2)

        # Hyperlink (simulated as underlined, blue text to mimic a hyperlink)
        hyperlink_paragraph = doc.add_paragraph()
        hyperlink_run = hyperlink_paragraph.add_run(TestFormattedTextLoading.FORMATTED_CONTENT["hyperlink_text"])
        hyperlink_run.font.underline = True
        hyperlink_run.font.color.rgb = RGBColor(0, 0, 255)  # Blue color for hyperlink simulation

        # List items
        doc.add_paragraph("List Items:")
        for item in TestFormattedTextLoading.FORMATTED_CONTENT["list_items"]:
            doc.add_paragraph(item, style="List Bullet")

        # Table
        table_data = TestFormattedTextLoading.FORMATTED_CONTENT["table_data"]
        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
        for i, row_data in enumerate(table_data):
            cells = table.rows[i].cells
            for j, cell_text in enumerate(row_data):
                cells[j].text = cell_text

        # Insert Image if available
        image_path = TestFormattedTextLoading.FORMATTED_CONTENT["image_path"]
        if os.path.exists(image_path):
            doc.add_paragraph("Image:")
            doc.add_picture(image_path, width=Pt(100))

        # Page break
        doc.add_page_break()

        doc.save(file_path)
        return file_path

    @staticmethod
    def create_formatted_pdf_file(folder_path, filename="formatted_test_file.pdf"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        pdf = TestFormattedTextLoading.PDF()
        pdf.add_page()

         # Add NotoSans TrueType Unicode font
        pdf.add_font("NotoSans", "", "tests/archived_files/Noto_Sans_JP/Noto_Sans/NotoSans-VariableFont_wdth,wght.ttf")
        pdf.add_font("NotoSans", "B", "tests/archived_files/Noto_Sans_JP/Noto_Sans/NotoSans-VariableFont_wdth,wght.ttf")
        pdf.add_font("NotoSans", "I", "tests/archived_files/Noto_Sans_JP/Noto_Sans/NotoSans-VariableFont_wdth,wght.ttf")
        pdf.set_font("NotoSans", "", size=12)


        # Add content with UTF-8 encoding support
        pdf.cell(0, 10, "Some text including a bullet point: •", new_x=XPos.LMARGIN, new_y=YPos.NEXT)  # Ensure it can handle UTF-8

        # Bold text
        pdf.set_font("NotoSans", style="B", size=12)
        pdf.cell(0, 10, "This text is bold.", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Italic text
        pdf.set_font("NotoSans", style="I", size=12)
        pdf.cell(0, 10, "This text is italic.", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Underlined text
        pdf.set_font("NotoSans", style="U", size=12)
        pdf.cell(0, 10, "This text is underlined.", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Colored text
        pdf.set_text_color(255, 0, 0)  # Red color
        pdf.cell(0, 10, "This text is in red color.", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Large text
        pdf.set_text_color(0, 0, 0)  # Reset color
        pdf.set_font("NotoSans", size=18)
        pdf.cell(0, 10, "This text is large.", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Centered text
        pdf.set_xy(10, pdf.get_y() + 10)
        pdf.cell(0, 10, "This text is centered.", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")

        # Headings
        pdf.set_font("NotoSans", "B", 16)
        pdf.cell(0, 10, "This is Heading 1", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("NotoSans", "B", 14)
        pdf.cell(0, 10, "This is Heading 2", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Hyperlink (simulated by underlining and coloring blue)
        pdf.set_font("NotoSans", style="U", size=12)
        pdf.set_text_color(0, 0, 255)  # Blue color for hyperlink simulation
        pdf.cell(0, 10, "OpenAI website", new_x=XPos.LMARGIN, new_y=YPos.NEXT, link="https://www.openai.com")

        # Reset text color
        pdf.set_text_color(0, 0, 0)

        # List items
        pdf.cell(0, 10, "List Items:", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        list_items = ["First item", "Second item", "Third item"]
        pdf.set_x(20)
        for item in list_items:
            pdf.cell(0, 10, f"• {item}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Table
        pdf.cell(0, 10, "Table Data:", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        table_data = [["Header 1", "Header 2"], ["Row 1, Cell 1", "Row 1, Cell 2"], ["Row 2, Cell 1", "Row 2, Cell 2"]]
        pdf.set_font("NotoSans", size=10)
        cell_width = 45
        cell_height = 10
        for row in table_data:
            for item in row:
                pdf.cell(cell_width, cell_height, item, border=1)
            pdf.ln(cell_height)

        # Image (if available)
        image_path = "tests/test_files/img.png"
        if os.path.exists(image_path):
            pdf.cell(0, 10, "Image:", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.image(image_path, x=10, y=pdf.get_y(), w=50)
            pdf.ln(60)

        # Save PDF
        pdf.output(file_path)
        return file_path

    @staticmethod
    def create_formatted_pptx_file(folder_path, filename="formatted_test_file.pptx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        ppt = Presentation()

        # Slide 1: Text Formatting
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = "Text Formatting Examples"

        # Bold, italic, underlined, colored text
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
        tf = textbox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "This text is bold."
        run.font.bold = True

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "This text is italic."
        run.font.italic = True

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "This text is underlined."
        run.font.underline = True

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "This text is in red color."
        run.font.color.rgb = PPTXColor(255, 0, 0)

        # Large text and centered text
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "This text is large."
        run.font.size = Pt(24)

        # Slide 2: Headings, Hyperlink, List, Table, Image
        slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide2.shapes.title.text = "More Formatting Examples"

        # Heading
        heading = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        heading.text_frame.add_paragraph().text = "This is Heading 1"
        heading.text_frame.add_paragraph().font.bold = True
        heading.text_frame.add_paragraph().font.size = Pt(20)

        # Hyperlink (simulated as blue, underlined text)
        hyperlink_paragraph = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1)).text_frame
        hyperlink = hyperlink_paragraph.add_paragraph()
        hyperlink.text = "OpenAI website"
        hyperlink.font.color.rgb = PPTXColor(0, 0, 255)
        hyperlink.font.underline = True

        # List items
        bullet_slide = slide2.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
        bullet_tf = bullet_slide.text_frame
        bullet_tf.add_paragraph().text = "List Items:"
        for item in ["First item", "Second item", "Third item"]:
            p = bullet_tf.add_paragraph()
            p.text = item
            p.level = 1  # Bullet level

        # Table
        table_slide = slide2.shapes.add_table(rows=3, cols=2, left=Inches(1), top=Inches(5), width=Inches(6), height=Inches(1.5)).table

        #set table headers
        table_slide.cell(0, 0).text = "Header 1"
        table_slide.cell(0, 1).text = "Header 2"

        table_data = [["Header 1", "Header 2"], ["Row 1, Cell 1", "Row 1, Cell 2"], ["Row 2, Cell 1", "Row 2, Cell 2"]]
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                table_slide.cell(row_idx, col_idx).text = cell_text

        # Image
        image_path = "tests/test_files/img.png"
        if os.path.exists(image_path):
            slide2.shapes.add_picture(image_path, Inches(1), Inches(6), width=Inches(2), height=Inches(2))

        # Save the presentation
        ppt.save(file_path)
        return file_path

    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for the test file
        self.folder_path = "tests/test_files"
        self.docx_file_path = self.create_formatted_docx_file(self.folder_path, "formatted_test_file.docx")
        self.pdf_file_path = self.create_formatted_pdf_file(self.folder_path, "formatted_test_file.pdf")
        self.pptx_file_path = self.create_formatted_pptx_file(self.folder_path, "formatted_test_file.pptx")

    def check_formatting(self, document_text, extracted_tables):
        # Verify bold and italic text
        assert self.FORMATTED_CONTENT["bold_text"] in document_text, "Bold text not preserved."
        assert self.FORMATTED_CONTENT["italic_text"] in document_text, "Italic text not preserved."

        # Verify list items
        for item in self.FORMATTED_CONTENT["list_items"]:
            assert item in document_text, f"List item '{item}' not preserved."

        # Verify table content
        expected_table_data = self.FORMATTED_CONTENT["table_data"]
        assert len(extracted_tables) > 0, "No tables found in the document."

        # Check the structure of the first table
        table = extracted_tables[0]  # Assuming we check the first table
        assert len(table) == len(expected_table_data), "Table row count does not match expected."
        for row_idx, expected_row in enumerate(expected_table_data):
            assert len(table[row_idx]) == len(expected_row), f"Row {row_idx} column count does not match expected."
            for col_idx, expected_value in enumerate(expected_row):
                cell_value = table[row_idx][col_idx]
                assert cell_value == expected_value, f"Expected '{expected_value}' in table cell ({row_idx}, {col_idx}), but found '{cell_value}'."

    def test_formatted_docx_file_loading(self):
        docx_loader = DOCXLoader(self.docx_file_path)
        docx_loader.validate_file()
        docx_document = docx_loader.load_file()
        assert docx_document is not None, "Failed to load the formatted DOCX document."

        # Extract text and tables from DOCX
        paragraphs = [p.text for p in docx_document.paragraphs]
        tables = [
            [  # Convert each table's rows and cells into lists
                [cell.text for cell in row.cells]
                for row in table.rows
            ]
            for table in docx_document.tables
        ]

        # Check formatting
        self.check_formatting(paragraphs, tables)

    def test_formatted_pdf_file_loading(self):
        pdf_loader = PDFLoader(self.pdf_file_path)
        pdf_loader.validate_file()
        pdf_document = pdf_loader.load_file()
        assert pdf_document is not None, "Failed to load the formatted PDF document."

        # Extract text and tables from PDF
        pdf_text_extractor = DataExtractor(pdf_loader)
        pdf_text = pdf_text_extractor.extract_text_from_pdf(pdf_document)
        tables = pdf_text_extractor.extract_tables()

        # Check formatting
        self.check_formatting(pdf_text, tables)

    def test_format_pptx_file_loading(self):
        pptx_loader = PPTLoader(self.pptx_file_path)
        pptx_loader.validate_file()
        pptx_document = pptx_loader.load_file()
        assert pptx_document is not None, "Failed to load the formatted PPTX document."

        # Extract text and tables from PPTX
        pptx_text_extractor = DataExtractor(pptx_loader)
        pptx_text = pptx_text_extractor.extract_text_from_pptx(pptx_document)
        tables = pptx_text_extractor.extract_tables()

        # Check formatting
        self.check_formatting(pptx_text, tables)