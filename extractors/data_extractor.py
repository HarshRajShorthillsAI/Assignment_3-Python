import fitz  # PyMuPDF for PDF handling
import camelot  # For PDF table extraction
import docx  # For DOCX handling
import pptx  # For PPTX handling
from PIL import Image as PILImage
import os
from io import BytesIO
import pandas as pd
from loaders.pdf_loader import PDFLoader
from loaders.docx_loader import DOCXLoader
from loaders.ppt_loader import PPTLoader

class DataExtractor:
    def __init__(self, loader):
        self.loader = loader
        self.file_path = loader.file_path

    def extract_text(self):
        loader = self.loader.load_file()
        if self.loader.get_fileType() == "pdf":
            return self.extract_text_from_pdf(loader)
        elif self.loader.get_fileType() == "docx":
            return self.extract_text_from_docx(loader)
        elif self.loader.get_fileType() == "pptx":
            return self.extract_text_from_pptx(loader)
        else:
            raise ValueError("Unsupported file format for text extraction.")

    def extract_text_from_pdf(self, reader):
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    
    def extract_text_from_docx(self, doc):
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    
    def extract_text_from_pptx(self, ppt):
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_images(self): #working in sql, defaulted to png in file.
        images = []
        if self.loader.get_fileType() == "pdf":
            images = self.extract_images_from_pdf()
        elif self.loader.get_fileType() == "docx":
            images = self.extract_images_from_docx()
        elif self.loader.get_fileType() == "pptx":
            images = self.extract_images_from_pptx()
        return images

    def extract_images_from_pdf(self):
        images = []
        pdf_document = fitz.open(self.file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                width, height = base_image["width"], base_image["height"]
                images.append({
                    "image_data": image_bytes,
                    "ext": image_ext,
                    "page": page_num + 1,
                    "dimensions": (width, height)
                })
        pdf_document.close()
        return images
    
    def extract_images_from_docx(self):
        images = []
        doc = self.loader.load_file()
        for rel in doc.part.rels:
            if doc.part.rels[rel].is_external == False and "image" in doc.part.rels[rel].target_ref:
                image_binary = doc.part.rels[rel].target_part.blob
                image_data = self.convert_to_pil(image_binary)
                images.append({"image_data": image_binary, "image": image_data})
        return images
    
    def extract_images_from_pptx(self):
        images = []
        ppt = self.loader.load_file()
        for slide_number, slide in enumerate(ppt.slides):
            for shape_number, shape in enumerate(slide.shapes):
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                    image_stream = shape.image.blob
                    image_ext = shape.image.ext
                    image = self.convert_to_pil(image_stream)
                    images.append({
                        'image': image,
                        'image_data': image_stream,
                        'ext': image_ext,
                        'slide_number': slide_number + 1,
                        'shape_number': shape_number + 1
                    })
        return images
    
    def convert_to_pil(self, image_stream):
        try:
            return PILImage.open(BytesIO(image_stream))
        except Exception as e:
            print("Failed to load image:", e)
            return None

    def extract_urls(self):
        urls = []
        if self.loader.get_fileType() == "pdf":
            urls = self.extract_urls_from_pdf()
        elif self.loader.get_fileType() == "docx":
            urls = self.extract_urls_from_docx()
        elif self.loader.get_fileType() == "pptx":
            urls = self.extract_urls_from_pptx()
        return urls
    
    # PDF, DOCX, and PPTX specific URL extraction methods follow similar pattern
    def extract_urls_from_pdf(self):
        urls = []
        pdf_document = fitz.open(self.file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            links = page.get_links()
            for link in links:
                if "uri" in link:
                    urls.append({
                        "url": link["uri"],
                        "page": page_num + 1,
                        "position": {
                            "x0": link["from"].x0,
                            "y0": link["from"].y0,
                            "x1": link["from"].x1,
                            "y1": link["from"].y1
                        }
                    })
        pdf_document.close()
        return urls
    
    def extract_urls_from_docx(self):
        urls = []
        # DOCX URL extraction
        doc = self.loader.load_file()

        for rel in doc.part.rels.values():
            # if "hyperlink" in rel.target_ref:
            urls.append({"url": rel.reltype})
        
        return urls

    def extract_urls_from_pptx(self):
        urls = []
        # PPTX URL extraction
        presentation = pptx.Presentation(self.file_path)

        for slide_num, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                # Check if the shape has a hyperlink
                if shape.has_text_frame and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink.address:
                                urls.append({
                                "url": run.hyperlink.address,
                                "slide": slide_num,
                                "text": run.text
                                })

                # Shapes like images or other elements may have hyperlinks as well
                if shape.has_text_frame is False and hasattr(shape, 'hyperlink') and shape.hyperlink and shape.hyperlink.target:
                    urls.append({
                        "url": shape.hyperlink.target,
                        "slide": slide_num,                        "shape_type": shape.shape_type
                    })
        return urls

    def extract_tables(self):
        if self.loader.get_fileType() == "pdf":
            return self.extract_tables_from_pdf()
        elif self.loader.get_fileType() == "docx":
            return self.extract_tables_from_docx()
        elif self.loader.get_fileType() == "pptx":
            return self.extract_tables_from_pptx()
        
    def extract_tables_from_pdf(self):
        tables = camelot.read_pdf(self.file_path, pages="all")
        return [table.df for table in tables]

    def extract_tables_from_docx(self):
        doc = self.loader.load_file()
        return [[
            [cell.text for cell in row.cells] for row in table.rows
        ] for table in doc.tables]

    def extract_tables_from_pptx(self):
        ppt = self.loader.load_file()
        tables = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    tables.append([
                        [cell.text for cell in row.cells] for row in shape.table.rows
                    ])
        return tables