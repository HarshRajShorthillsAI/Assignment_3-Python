import os
import pytest
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import fitz

from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
from loaders.ppt_loader import PPTLoader

class TestUTF8FileLoading:
    @staticmethod
    def create_utf8_encoded_pdf_file(folder_path, filename="utf8_test_file.pdf"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, filename)

        try:
            # Initialize FPDF
            pdf = FPDF()
            pdf.add_page()

            # Add a TTF font that supports UTF-8 characters
            font_path = "tests/archived_files/Noto_Sans_JP/Noto_Sans_JP/NotoSansJP-VariableFont_wght.ttf"
            pdf.add_font("NotoSansJP", "", font_path, uni=True)
            pdf.set_font("NotoSansJP", size=12)

            # UTF-8 encoded content with special characters
            utf8_content = "This is a UTF-8 encoded file with special characters: ñ, ü, é, ö, ç, こんにちは"

            # Add content to PDF
            pdf.multi_cell(0, 10, utf8_content)

            # Save the PDF
            pdf.output(file_path)
            return file_path

        except FileNotFoundError:
            print(f"Font file not found at {font_path}. Ensure the path is correct.")
        except Exception as e:
            print(f"An error occurred: {e}")

        return None

    @staticmethod
    def create_utf8_encoded_docx_file(folder_path, filename="utf8_test_file.docx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        content = "This is a UTF-8 encoded file with special characters: ñ, ü, é, ö, ç, こんにちは"
        doc = Document()
        doc.add_paragraph(content)
        doc.save(file_path)

        return file_path

    @staticmethod
    def create_utf8_encoded_pptx_file(folder_path, filename="utf8_test_file.pptx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        content = "This is a UTF-8 encoded file with special characters: ñ, ü, é, ö, ç, こんにちは"
        ppt = Presentation()
        slide_layout = ppt.slide_layouts[1]  # Title and Content layout
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = "UTF-8 Test Slide"
        slide.placeholders[1].text = content
        ppt.save(file_path)

        return file_path

    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for test files
        self.folder_path = "tests/test_files"
        self.pdf_file_path = self.create_utf8_encoded_pdf_file(self.folder_path, "utf8_test_file.pdf")
        self.docx_file_path = self.create_utf8_encoded_docx_file(self.folder_path, "utf8_test_file.docx")
        self.pptx_file_path = self.create_utf8_encoded_pptx_file(self.folder_path, "utf8_test_file.pptx")

    def test_utf8_encoded_pdf_file_loading(self):
        # Initialize the PDFLoader with the path to the test PDF file
        pdf_loader = PDFLoader(self.pdf_file_path)

        # Use PDFLoader's load_file method to open and read the PDF
        pdf_document = pdf_loader.load_file()
        assert pdf_document is not None, "Failed to load the UTF-8 encoded PDF document."

    def test_utf8_encoded_docx_file_loading(self):
        #Initialize the DOCXLoader with the path to the test DOCX file
        docx_loader = DOCXLoader(self.docx_file_path)
        
        #Use DOCXLoader's load_file method to load the docx
        docx_document = docx_loader.load_file()
        
        assert docx_document is not None, "Failed to load the UTF-8 encoded DOCX document"

    def test_utf8_encoded_pptx_file_loading(self):
        #Initialize the PPTXLoader with the path to the test PPTX file
        pptx_loader = PPTLoader(self.pptx_file_path)

        #Use PPTLoader's load_file method to load the pptx
        pptx_document = pptx_loader.load_file()

        assert pptx_document is not None, "Failed to load the UTF-8 encoded PPTX document."