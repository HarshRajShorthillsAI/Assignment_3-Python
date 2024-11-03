import os
import pytest
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches  # Ensure Pt is imported
from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
from loaders.ppt_loader import PPTLoader
import io

class TestImageOnlyDocumentLoader:
    image_paths = ['tests/test_files/img.png']

    @staticmethod
    def create_image_only_pdf_file(folder_path, filename="image_only_file.pdf"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        
        for image_path in TestImageOnlyDocumentLoader.image_paths:
            pdf.image(image_path, x=10, w=190)  # Adjust width (w) as needed
            pdf.ln(10)  # Add space after each image

        pdf.output(file_path)
        return file_path  # Return the file path

    @staticmethod
    def create_image_only_docx_file(folder_path, filename="image_only_file.docx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)
        doc = Document()

        for image_path in TestImageOnlyDocumentLoader.image_paths:
            doc.add_picture(image_path, width=Inches(4))  # Adjust the width as needed
            doc.add_paragraph()  # Add a paragraph to provide spacing between images

        doc.save(file_path)
        return file_path  # Return the file path

    @staticmethod
    def create_image_only_pptx_file(folder_path, filename="image_only_file.pptx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)
        prs = Presentation()

        for image_path in TestImageOnlyDocumentLoader.image_paths:
            slide_layout = prs.slide_layouts[5]  # Using a blank layout
            slide = prs.slides.add_slide(slide_layout)
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            slide.shapes.add_picture(image_path, left, top, width, height)

        prs.save(file_path)
        return file_path  # Return the file path

    @pytest.fixture(autouse=True)
    def setup(self):
        self.folder_path = "tests/test_files"
        self.pdf_file_path = self.create_image_only_pdf_file(self.folder_path, "image_only_file.pdf")
        self.docx_file_path = self.create_image_only_docx_file(self.folder_path, "image_only_file.docx")
        self.pptx_file_path = self.create_image_only_pptx_file(self.folder_path, "image_only_file.pptx")

    def save_images_to_docx(self, image_data_list, output_docx_path):
        doc = Document()  # Create a new DOCX document
        
        for img_data in image_data_list:
            # Save the image data to a BytesIO object
            image_stream = io.BytesIO(img_data)
            doc.add_picture(image_stream)  # Add the image to the document
            doc.add_paragraph()  # Add a paragraph for spacing
        
        doc.save(output_docx_path)  # Save the DOCX file

    def test_images_only_pdf_file_loading(self):
        pdf_loader = PDFLoader(self.pdf_file_path)
        pdf_loader.validate_file()
        pdf_document = pdf_loader.load_file()

        assert pdf_document is not None, "Failed to load the image-only PDF document."

        images = []

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            image_list = page.get_images(full=True)  # Get all images on the page

            for img_index, img in enumerate(image_list):
                xref = img[0]  # Get the xref of the image
                base_image = pdf_document.extract_image(xref)
                image_data = base_image["image"]  # The binary data of the image
                images.append(image_data)

        assert len(images) > 0, "No images found in the PDF document."


        os.makedirs("extracted_data/image_only", exist_ok=True)
        # Save the extracted images to a new DOCX file
        output_docx_path = os.path.join("extracted_data/image_only", "extracted_images_from_pdf.docx")

        self.save_images_to_docx(images, output_docx_path)

    def test_image_only_docx_file_loading(self):
        docx_loader = DOCXLoader(self.docx_file_path)
        docx_loader.validate_file()
        docx_document = docx_loader.load_file()
        assert docx_document is not None, "Failed to load the text only DOCX document"

        # Extract images from the DOCX
        extracted_images = []
        for rel in docx_document.part.rels.values():
            if "image" in rel.reltype:
                extracted_images.append(rel.target_part.blob)  # Get image data


        assert len(extracted_images) > 0, "No images found in the DOCX document."

        os.makedirs("extracted_data/image_only", exist_ok=True)
        # Save the extracted images to a new DOCX file
        output_docx_path = os.path.join("extracted_data/image_only", "extracted_images_from_docx.docx")

        self.save_images_to_docx(extracted_images, output_docx_path)

    def test_image_only_pptx_file_loading(self):
        pptx_loader = PPTLoader(self.pptx_file_path)
        pptx_loader.validate_file()
        pptx_document = pptx_loader.load_file()
        assert pptx_document is not None, "Failed to load the text only PPTX document."

        # Extract images from the PPTX
        extracted_images = []
        for slide in pptx_document.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    extracted_images.append(shape.image.blob)  # Get image data

        assert len(extracted_images) > 0, "No images found in the PPTX document."

        os.makedirs("extracted_data/image_only", exist_ok=True)
        # Save the extracted images to a new DOCX file
        output_docx_path = os.path.join("extracted_data/image_only", "extracted_images_from_pptx.docx")

        # Store the extracted images into a DOCX file
        self.save_images_to_docx(extracted_images, output_docx_path)