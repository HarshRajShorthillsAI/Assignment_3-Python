from docx import Document
import fitz
from pptx import Presentation  # PyMuPDF
from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
import pytest
import os

from loaders.ppt_loader import PPTLoader

class TestTableOnlyDocumentLoader:
    # Sample data for tables
    tables = {
        "Employee Info": [
            ["ID", "Name", "Age", "Position", "Department"],
            [101, "Alice Johnson", 28, "Software Engineer", "IT"],
            [102, "Bob Smith", 35, "Product Manager", "Marketing"],
            [103, "Carol Williams", 40, "HR Manager", "Human Resources"],
            [104, "David Brown", 32, "Data Analyst", "IT"],
            [105, "Eva Green", 29, "Designer", "Design"]
        ],
        "Salary Info": [
            ["ID", "Name", "Base Salary", "Bonus", "Total Salary"],
            [101, "Alice Johnson", "$70,000", "$5,000", "$75,000"],
            [102, "Bob Smith", "$90,000", "$10,000", "$100,000"],
            [103, "Carol Williams", "$65,000", "$7,000", "$72,000"],
            [104, "David Brown", "$80,000", "$6,000", "$86,000"],
            [105, "Eva Green", "$55,000", "$4,000", "$59,000"]
        ],
        "Project Assignments": [
            ["ID", "Name", "Project", "Role", "Duration (months)"],
            [101, "Alice Johnson", "Project A", "Lead Developer", 12],
            [102, "Bob Smith", "Project B", "Product Owner", 8],
            [103, "Carol Williams", "Project C", "HR Liaison", 6],
            [104, "David Brown", "Project A", "Data Specialist", 10],
            [105, "Eva Green", "Project D", "Lead Designer", 9]
        ],
        "Attendance Records": [
            ["ID", "Name", "Month", "Days Present", "Days Absent"],
            [101, "Alice Johnson", "January", 20, 1],
            [102, "Bob Smith", "January", 19, 2],
            [103, "Carol Williams", "January", 21, 0],
            [104, "David Brown", "January", 18, 3],
            [105, "Eva Green", "January", 22, 0]
        ],
        "Training Programs": [
            ["ID", "Name", "Program", "Completion Date", "Status"],
            [101, "Alice Johnson", "Advanced Python", "2023-04-15", "Completed"],
            [102, "Bob Smith", "Leadership Skills", "2023-05-10", "Completed"],
            [103, "Carol Williams", "HR Management", "2023-06-25", "Completed"],
            [104, "David Brown", "Data Analysis", "2023-07-05", "In Progress"],
            [105, "Eva Green", "Design Fundamentals", "2023-08-12", "Completed"]
        ],
        "Performance Reviews": [
            ["ID", "Name", "Review Period", "Score", "Comments"],
            [101, "Alice Johnson", "2023-Q1", 4.5, "Excellent technical skills"],
            [102, "Bob Smith", "2023-Q1", 4.2, "Strong leadership"],
            [103, "Carol Williams", "2023-Q1", 4.0, "Effective team player"],
            [104, "David Brown", "2023-Q1", 3.8, "Good analytical skills"],
            [105, "Eva Green", "2023-Q1", 4.6, "Creative and detail-oriented"]
        ],
        "Employee Benefits": [
            ["ID", "Name", "Health Insurance", "Retirement Plan", "Vacation Days"],
            [101, "Alice Johnson", "Plan A", "5% Match", 20],
            [102, "Bob Smith", "Plan B", "6% Match", 18],
            [103, "Carol Williams", "Plan C", "4% Match", 25],
            [104, "David Brown", "Plan A", "5% Match", 15],
            [105, "Eva Green", "Plan B", "6% Match", 20]
        ],
        "Promotion History": [
            ["ID", "Name", "Previous Position", "New Position", "Date of Promotion"],
            [101, "Alice Johnson", "Junior Developer", "Software Engineer", "2022-02-20"],
            [102, "Bob Smith", "Marketing Associate", "Product Manager", "2021-06-15"],
            [103, "Carol Williams", "HR Assistant", "HR Manager", "2022-01-10"],
            [104, "David Brown", "Intern", "Data Analyst", "2021-08-25"],
            [105, "Eva Green", "Graphic Designer", "Designer", "2023-03-05"]
        ],
        "Overtime Records": [
            ["ID", "Name", "Month", "Overtime Hours", "Overtime Pay"],
            [101, "Alice Johnson", "January", 5, "$300"],
            [102, "Bob Smith", "January", 8, "$480"],
            [103, "Carol Williams", "January", 4, "$240"],
            [104, "David Brown", "January", 10, "$600"],
            [105, "Eva Green", "January", 6, "$360"]
        ],
        "Leave Requests": [
            ["ID", "Name", "Leave Type", "Start Date", "Status"],
            [101, "Alice Johnson", "Sick Leave", "2023-02-01", "Approved"],
            [102, "Bob Smith", "Vacation", "2023-02-10", "Pending"],
            [103, "Carol Williams", "Maternity Leave", "2023-03-01", "Approved"],
            [104, "David Brown", "Personal Leave", "2023-04-15", "Rejected"],
            [105, "Eva Green", "Vacation", "2023-05-10", "Approved"]
        ]
    }

    @staticmethod
    def create_table_only_pdf_file(folder_path, file_name="table_only_pdf_file.pdf"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, file_name)

        try:

            # Initialize the PDF document
            doc = fitz.open()

            # Styling and layout parameters
            font_size = 10
            cell_widths = [50, 120, 100, 100, 100, 100]  # Adjust widths to fit table (6 columns max)
            row_height = 20
            left_margin = 72  # 1 inch from the left
            top_margin = 72   # 1 inch from the top
            table_spacing = 20  # Space between tables
            page_height = 792  # Default height for letter size
            content_height = page_height - top_margin  # Max content height

            # Create a new page
            page = doc.new_page()  

            # Calculate initial y position for the first table
            y = top_margin

            # Generate each table on the same page
            for table_title, data in TestTableOnlyDocumentLoader.tables.items():
                # Check if there's enough space for the title and table
                if y + 20 + (len(data) * row_height) + table_spacing > page_height:
                    # Create a new page if not enough space
                    page = doc.new_page()
                    y = top_margin  # Reset y position for the new page

                # Draw the table title
                page.insert_text((left_margin, y), table_title, fontname="helv", fontsize=12, color=(0, 0, 0))
                
                # Move y position down to start drawing the table rows
                y += 20  # Adjust this to control spacing after the title
                
                # Draw each row and column
                for row in data:
                    x = left_margin
                    
                    # Draw each cell in the row
                    for i, item in enumerate(row):
                        # Check if the current index is within bounds of cell_widths
                        if i >= len(cell_widths):
                            print(f"Warning: Column index {i} exceeds defined cell widths for table '{table_title}'.")
                            break
                        
                        # Calculate cell rectangle
                        rect = fitz.Rect(x, y, x + cell_widths[i], y + row_height)
                        
                        # Draw cell border
                        page.draw_rect(rect, color=(0, 0, 0), width=0.5)

                        # Insert text into the cell, centered
                        page.insert_textbox(
                            rect,                      # Rectangle for text
                            str(item),                 # Text content
                            fontname="helv",           # Font name
                            fontsize=font_size,        # Font size
                            color=(0, 0, 0),           # Text color
                            align=fitz.TEXT_ALIGN_CENTER  # Center-align the text
                        )
                        
                        # Move to the next cell
                        x += cell_widths[i]

                    # Move y position down for the next row
                    y += row_height
                
                # Move y position down for the next table
                y += table_spacing

            # Save the PDF to the specified file path
            doc.save(file_path)
            doc.close()
            return file_path
        except Exception as e:
            print(f"An error occured: {e}")

        return None

    @staticmethod
    def create_table_only_docx_file(folder_path, file_name="table_only_docx_file.docx"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)

        # Define the full path for the DOCX file
        file_path = os.path.join(folder_path, file_name)

        # Create a new Document
        doc = Document()

        # Generate each table in the DOCX file
        for table_title, data in TestTableOnlyDocumentLoader.tables.items():
            # Add the table title
            doc.add_heading(table_title, level=1)

            # Create a new table
            table = doc.add_table(rows=0, cols=len(data[0]))  # Create table with columns equal to the first row

            # Add rows to the table
            for row in data:
                cells = table.add_row().cells
                for i, item in enumerate(row):
                    cells[i].text = str(item)  # Insert text into each cell

        # Save the document
        doc.save(file_path)
        return file_path

    @staticmethod
    def create_table_only_pptx_file(folder_path, file_name="table_only_pptx_file.pptx"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)

        # Define the full path for the PPTX file
        file_path = os.path.join(folder_path, file_name)

        # Create a new Presentation
        prs = Presentation()

        # Generate each table in the PPTX file
        for table_title, data in TestTableOnlyDocumentLoader.tables.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
            title = slide.shapes.title
            title.text = table_title

            # Define table dimensions and positions
            rows, cols = len(data), len(data[0])
            left = top = width = height = 0  # Define dimensions for the table here

            # Create a table
            table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

            # Set column widths if necessary
            for i in range(cols):
                table.columns[i].width = int(width / cols)  # Adjust as needed

            # Add header row
            for i, item in enumerate(data[0]):
                table.cell(0, i).text = str(item)

            # Add data rows
            for row_idx in range(1, rows):
                for col_idx in range(cols):
                    table.cell(row_idx, col_idx).text = str(data[row_idx][col_idx])

        # Save the presentation
        prs.save(file_path)
        return file_path


    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for test files
        self.folder_path = "tests/test_files"
        self.pdf_file_path = self.create_table_only_pdf_file(self.folder_path, "table_only_file.pdf")
        self.docx_file_path = self.create_table_only_docx_file(self.folder_path, "table_only_file.docx")
        self.pptx_file_path = self.create_table_only_pptx_file(self.folder_path, "table_only_file.pptx")

    def test_table_only_pdf_file_loading(self):
        # Initialize the PDFLoader with the path to the test PDF file
        pdf_loader = PDFLoader(self.pdf_file_path)
        pdf_loader.validate_file()

        # Use PDFLoader's load_file method to open and read the PDF
        pdf_document = pdf_loader.load_file()
        assert pdf_document is not None, "Failed to load the table-only PDF document."

        # Extract text and check for tables
        extracted_tables = []
        for page in pdf_document:
            text = page.get_text("text")  # Extract text as plain text
            lines = text.splitlines()

            table = []
            for line in lines:
                # Assuming space or tab as the delimiter, customize as necessary
                columns = line.split('\t')
                if columns:  # Only add non-empty lines
                    table.append(columns)

            if table:  # Append the table if it's not empty
                extracted_tables.append(table)

        # Check if any tables were extracted
        assert len(extracted_tables) > 0, "No tables found in the PDF document."

        # Display the loaded table data
        print("Loaded Table Data:\n")
        for table in extracted_tables:
            for row in table:
                print(row, sep="\t")
            print('')

    def test_table_only_docx_file_loading(self):
        # Load the DOCX file and check for tables
        docx_loader = DOCXLoader(self.docx_file_path)
        doc = docx_loader.validate_file()
        doc = docx_loader.load_file()

        extracted_tables = []
        for table in doc.tables:
            extracted_table = []
            for row in table.rows:
                extracted_row = [cell.text for cell in row.cells]
                extracted_table.append(extracted_row)
            extracted_tables.append(extracted_table)

        # Check if any tables were extracted
        assert len(extracted_tables) > 0, "No tables found in the DOCX document."

        # Display the loaded table data
        print("Loaded Table Data from DOCX:\n")
        for table in extracted_tables:
            for row in table:
                print(row, sep="\t")
            print('')

    def test_table_only_pptx_file_loading(self):
        # Load the PPTX file and check for tables
        pptx_loader = PPTLoader(self.pptx_file_path)
        doc = pptx_loader.validate_file()
        doc = pptx_loader.load_file()

        extracted_tables = []
        for slide in doc.slides:
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    extracted_table = []
                    table = shape.table
                    for row in table.rows:
                        extracted_row = [cell.text for cell in row.cells]
                        extracted_table.append(extracted_row)
                    extracted_tables.append(extracted_table)

        # Check if any tables were extracted
        assert len(extracted_tables) > 0, "No tables found in the PPTX document."

        # Display the loaded table data
        print("Loaded Table Data from PPTX:\n")
        for table in extracted_tables:
            for row in table:
                print(row, sep="\t")
            print('')