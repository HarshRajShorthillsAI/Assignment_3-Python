import os
import pytest
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import fitz
from pptx.util import Inches, Pt  # Ensure Pt is imported

from pptx.enum.text import MSO_AUTO_SIZE

from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
from loaders.ppt_loader import PPTLoader

class TestTextOnlyDocumentLoader:# UTF-8 encoded content with special characters
    text_content = ''' The Rise of Large Language Models: Transforming Communication and Technology

    Introduction

    In recent years, the field of artificial intelligence (AI) has witnessed an extraordinary evolution, particularly with the advent of large language models (LLMs). These models have transformed the way we interact with technology, reshaping applications in numerous sectors, including education, healthcare, customer service, and creative industries. LLMs are designed to understand, generate, and manipulate human language, enabling machines to communicate in ways that were previously unimaginable. This essay explores the development, functionalities, implications, and future prospects of LLMs, providing a comprehensive overview of their significance in contemporary society.

    What Are Large Language Models?

    Large language models are a subset of AI that utilizes deep learning techniques to process and generate human-like text. They are trained on vast datasets containing text from diverse sources, such as books, articles, websites, and social media. By analyzing patterns in language, LLMs learn grammar, context, facts, and even some level of reasoning. The most notable examples of LLMs include OpenAI's GPT-3, Google's BERT, and Facebook's RoBERTa. These models can perform various tasks, such as text completion, translation, summarization, and question-answering, often achieving results that are indistinguishable from those produced by humans.

    The architecture of LLMs typically relies on neural networks, particularly the transformer model, which was introduced in the paper "Attention is All You Need" by Vaswani et al. in 2017. The transformer architecture enables the model to weigh the importance of different words in a sentence and capture long-range dependencies, making it particularly effective for understanding and generating language.

    Development and Training

    The training of LLMs involves several stages, including data collection, pre-processing, model design, and fine-tuning. First, vast amounts of textual data are collected from various sources to create a comprehensive training set. This data is then cleaned and processed to remove noise and irrelevant information. Next, the model architecture is designed, with parameters adjusted to optimize performance.

    Training LLMs requires significant computational resources and time. The process typically involves using powerful GPUs or TPUs in data centers. As the models learn from the training data, they adjust their internal parameters to minimize errors in predicting the next word in a sentence. This process can take days or even weeks, depending on the model's size and the dataset's complexity. After training, models can be fine-tuned on specific tasks or domains to enhance their performance in particular applications.

    Applications of Large Language Models

    The versatility of LLMs has led to their widespread adoption across various industries. In customer service, chatbots powered by LLMs can provide instant responses to user inquiries, improving efficiency and user satisfaction. In education, LLMs can serve as personalized tutors, adapting their teaching methods to individual learning styles and needs. They can also assist in creating educational content, generating quizzes, and summarizing lessons.

    In the healthcare sector, LLMs can analyze medical literature, aiding in research and diagnosis. They can assist healthcare professionals by providing insights based on the latest findings and suggesting potential treatment options. Furthermore, LLMs can help in drafting patient reports and managing administrative tasks, freeing up valuable time for medical staff.

    The creative industry has also embraced LLMs, using them to generate stories, poetry, and even code. Writers can leverage these models to overcome writer's block, brainstorm ideas, and refine their narratives. In journalism, LLMs can assist in drafting articles and summarizing news, allowing reporters to focus on more in-depth investigations.

    Ethical Considerations and Challenges

    Despite their potential, the deployment of LLMs raises ethical concerns and challenges. One major issue is the potential for bias in the models. Since LLMs are trained on data collected from the internet, they can inadvertently learn and propagate harmful stereotypes or biased language. For instance, if a model is trained on text that reflects societal biases, it may generate outputs that reinforce these biases, leading to discrimination or misinformation.

    Another concern is the environmental impact of training large models. The computational resources required for training LLMs consume significant energy, raising questions about sustainability in AI development. Researchers are exploring methods to create more efficient models and reduce the carbon footprint associated with training.

    Moreover, the misuse of LLMs poses a significant risk. These models can generate convincing fake news, impersonate individuals, or produce misleading information. The ease of generating realistic text may lead to the spread of misinformation, impacting public opinion and societal trust.

    The Future of Large Language Models

    As technology continues to advance, the future of LLMs holds both promise and challenges. Researchers are actively working to improve the capabilities of LLMs, focusing on enhancing their understanding of context, reasoning, and common sense. Future models may incorporate multimodal capabilities, allowing them to process not just text but also images, audio, and video, leading to more holistic AI systems.

    In addition to improving the models, there is a growing emphasis on responsible AI development. Initiatives aimed at addressing bias, promoting transparency, and ensuring ethical usage are becoming increasingly important. Collaboration between researchers, developers, and policymakers will be crucial in establishing guidelines and regulations for the responsible deployment of LLMs.

    Conclusion

    Large language models have revolutionized the way we interact with technology, providing new avenues for communication, creativity, and efficiency. As these models continue to evolve, their impact on society will grow, influencing various sectors and reshaping our understanding of AI. However, the ethical implications and challenges associated with LLMs must be addressed to ensure their responsible and beneficial use. By fostering collaboration and promoting responsible AI practices, we can harness the power of LLMs to create a future where technology enhances human capabilities and contributes to societal well-being. The journey of LLMs is just beginning, and their potential to transform our world is both exciting and daunting.
    '''

    @staticmethod
    def create_text_only_pdf_file(folder_path, filename="text_only_file.pdf"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, filename)
        try:
            # Initialize FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", "B",8)

            # Add content to PDF
            pdf.multi_cell(0, 10, TestTextOnlyDocumentLoader.text_content)

            # Save the PDF
            pdf.output(file_path)
            return file_path

        except Exception as e:
            print(f"An error occurred: {e}")

        return None

    @staticmethod
    def create_text_only_docx_file(folder_path, filename="text_only_file.docx"):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)
        doc = Document()
        doc.add_paragraph(TestTextOnlyDocumentLoader.text_content)
        doc.save(file_path)

        return file_path

    from pptx.util import Inches, Pt

    @staticmethod #need correction for content writing on pptx
    def create_text_only_pptx_file(folder_path, file_name):
        # Create a presentation
        prs = Presentation()

        # Define max width and height for the text box to avoid overflow
        textbox_width = Inches(8)  # Width constraint within slide
        textbox_height = Inches(5)  # Height constraint within slide

        # Split content into 50-word chunks
        words = TestTextOnlyDocumentLoader.text_content.split()
        chunks = [" ".join(words[i:i + 100]) for i in range(0, len(words), 100)]

        # Add a slide for each chunk
        for chunk in chunks:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
            content_shape = slide.shapes.add_textbox(Inches(1), Inches(1), textbox_width, textbox_height)
            text_frame = content_shape.text_frame

            # Set text and ensure wrapping
            p = text_frame.add_paragraph()
            p.text = chunk
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Prevents overflow by not resizing the shape
            text_frame.word_wrap = True  # Ensures words wrap within text box dimensions

            # Optional styling: adjust font size as needed
            p.font.size = Pt(18)  # Adjust font size if needed

        # Save the presentation
        pptx_path = f"{folder_path}/{file_name}"
        prs.save(pptx_path)
        return pptx_path


    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for test files
        self.folder_path = "tests/test_files"
        self.pdf_file_path = self.create_text_only_pdf_file(self.folder_path, "text_only_file.pdf")
        self.docx_file_path = self.create_text_only_docx_file(self.folder_path, "text_only_file.docx")
        self.pptx_file_path = self.create_text_only_pptx_file(self.folder_path, "text_only_file.pptx")

    def test_text_only_pdf_file_loading(self):
        # Initialize the PDFLoader with the path to the test PDF file
        pdf_loader = PDFLoader(self.pdf_file_path)
        pdf_loader.validate_file()

        # Use PDFLoader's load_file method to open and read the PDF
        pdf_document = pdf_loader.load_file()
        assert pdf_document is not None, "Failed to load the text only PDF document."

        # Extract text from the PDF document
        extracted_text = ""
        for page in pdf_document:
            extracted_text += page.get_text()

        # Display the loaded PDF text
        print("Loaded PDF Text Content:\n")
        print(extracted_text)  # Display the entire extracted text on terminal

    def test_text_only_docx_file_loading(self):
        #Initialize the DOCXLoader with the path to the test DOCX file
        docx_loader = DOCXLoader(self.docx_file_path)
        docx_loader.validate_file()
        #Use DOCXLoader's load_file method to load the docx
        docx_document = docx_loader.load_file()
        
        assert docx_document is not None, "Failed to load the text only DOCX document"

        # Extract text from the Document instance
        extracted_text = []
        for paragraph in docx_document.paragraphs:
            extracted_text.append(paragraph.text)

        # Combine the paragraphs into a single string
        full_text = "\n".join(extracted_text)

        # Display the loaded DOCX text
        print("Loaded DOCX Text Content:\n")
        print(full_text)  # Display the entire extracted text

    def test_text_only_pptx_file_loading(self):
        #Initialize the PPTXLoader with the path to the test PPTX file
        pptx_loader = PPTLoader(self.pptx_file_path)
        pptx_loader.validate_file()
        #Use PPTLoader's load_file method to load the pptx
        pptx_document = pptx_loader.load_file()

        assert pptx_document is not None, "Failed to load the text only PPTX document."

        # Extract text from the presentation
        extracted_text = []
        for slide in pptx_document.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Ensure the shape has text
                    extracted_text.append(shape.text)

        # Combine the text from all slides into a single string
        full_text = "\n".join(extracted_text)

        # Display the loaded PPTX text
        print("Loaded PPTX Text Content:\n")
        print(full_text)  # Display the entire extracted text