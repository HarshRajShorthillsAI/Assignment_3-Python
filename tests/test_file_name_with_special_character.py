import os
import pytest
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from loaders.pdf_loader import PDFLoader  # Assuming you have a PDFLoader class for loading PDFs
from loaders.docx_loader import DOCXLoader  # Assuming a DOCXLoader class exists
from loaders.ppt_loader import PPTLoader  # Assuming a PPTXLoader class exists

class TestSpecialCharacterFileLoading:
    @staticmethod
    def create_dummy_pdf_with_random_text(folder_path, filename="file@name.pdf", num_paragraphs=3):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="Special Character File Test - PDF", ln=True, align='C')
        for _ in range(num_paragraphs):
            pdf.multi_cell(0, 10, txt="Random text for PDF content simulation.")
            pdf.ln(5)

        pdf.output(file_path)
        return file_path

    @staticmethod
    def create_dummy_docx_with_random_text(folder_path, filename="file@name.docx", num_paragraphs=3):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        doc = Document()
        doc.add_heading("Special Character File Test - DOCX", level=1)
        for _ in range(num_paragraphs):
            doc.add_paragraph("Random text for DOCX content simulation.")
        
        doc.save(file_path)
        return file_path

    @staticmethod
    def create_dummy_pptx_with_random_text(folder_path, filename="file@name.pptx", num_slides=3):
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        ppt = Presentation()
        for _ in range(num_slides):
            slide_layout = ppt.slide_layouts[1]  # Title and Content layout
            slide = ppt.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Special Character File Test - PPTX"
            content.text = "Random text for PPTX content simulation."

        ppt.save(file_path)
        return file_path

    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for test files
        self.folder_path = "tests/test_files"
        self.pdf_file = self.create_dummy_pdf_with_random_text(self.folder_path, "file@name.pdf")
        self.docx_file = self.create_dummy_docx_with_random_text(self.folder_path, "file@name.docx")
        self.pptx_file = self.create_dummy_pptx_with_random_text(self.folder_path, "file@name.pptx")

        yield

        # Clean up created files after the test
        # os.remove(self.pdf_file)
        # os.remove(self.docx_file)
        # os.remove(self.pptx_file)

    def test_special_character_pdf_loading(self):
        pdf_loader = PDFLoader(self.pdf_file)
        docs = pdf_loader.load_file()
        assert docs, "Failed to load the PDF file with special characters in the filename."

    def test_special_character_docx_loading(self):
        docx_loader = DOCXLoader(self.docx_file)
        docs = docx_loader.load_file()
        assert docs, "Failed to load the DOCX file with special characters in the filename."

    def test_special_character_pptx_loading(self):
        pptx_loader = PPTLoader(self.pptx_file)
        docs = pptx_loader.load_file()
        assert docs, "Failed to load the PPTX file with special characters in the filename."