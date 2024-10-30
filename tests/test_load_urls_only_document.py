from fpdf import FPDF
import os

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTXColor
from pptx import Presentation
from docx.shared import RGBColor

import docx
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

import pytest

from loaders.docx_loader import DOCXLoader
from loaders.pdf_loader import PDFLoader
from loaders.ppt_loader import PPTLoader

class TestURLOnlyDocumentLoader:
    # Define the data with URLs, pages, and positions
    data = [
        {'url': 'https://openreview.net/forum?id=R8sQPpGCv0', 'page': 37, 'position': {'x0': 163.82, 'y0': 595.89, 'x1': 289.67, 'y1': 604.99}},
        {'url': 'https://openreview.net/forum?id=R8sQPpGCv0', 'page': 37, 'position': {'x0': 58.86, 'y0': 605.35, 'x1': 213.59, 'y1': 614.45}},
        {'url': 'https://openreview.net/forum?id=R8sQPpGCv0', 'page': 37, 'position': {'x0': 76.79, 'y0': 624.13, 'x1': 256.58, 'y1': 633.43}},
        {'url': 'https://github.com/OpenBMB/BMTrain', 'page': 37, 'position': {'x0': 327.85, 'y0': 330.88, 'x1': 470.05, 'y1': 339.98}},
        {'url': 'https://github.com/OpenBMB/BMTrain', 'page': 37, 'position': {'x0': 345.79, 'y0': 340.19, 'x1': 491.70, 'y1': 349.49}},
        # Additional data is truncated for brevity
    ]

    @staticmethod
    def create_url_only_pdf_file(folder_path, filename="url_only_file.pdf"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, filename)
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
    
        for item in TestURLOnlyDocumentLoader.data:
            url = item['url']
            pdf.cell(0, 10, txt=url, ln=True, link=url)
        
        pdf.output(file_path)
        return file_path


    @staticmethod
    def create_url_only_docx_file(folder_path, filename="url_only_file.docx"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, filename)
        
        doc = Document()
    
        for item in TestURLOnlyDocumentLoader.data:
            paragraph = doc.add_paragraph()
            hyperlink = TestURLOnlyDocumentLoader.create_hyperlink(doc, item['url'])
            paragraph._p.append(hyperlink)
        
        doc.save(file_path)
        return file_path
    
    @staticmethod
    def create_hyperlink(doc, url):
        # Create hyperlink element
        part = doc.part
        r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        
        # Create hyperlink text run
        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        color = OxmlElement("w:color")
        color.set(qn("w:val"), "0000FF")  # Hex color for blue
        rPr.append(color)
        new_run.append(rPr)
        
        # Add link text
        text = OxmlElement("w:t")
        text.text = url
        new_run.append(text)
        hyperlink.append(new_run)
        
        return hyperlink

    
    @staticmethod
    def create_url_only_pptx_file(folder_path, filename="url_only_file.pptx"):
        # Ensure the folder exists
        os.makedirs(folder_path, exist_ok=True)
        
        # Define the full path for the PDF file
        file_path = os.path.join(folder_path, filename)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
        
        for item in TestURLOnlyDocumentLoader.data:
            url = item['url']
            x = Inches(1)  # Example x position, replace with item['position']['x0'] if needed
            y = Inches(1)  # Example y position, replace with item['position']['y0'] if needed
            
            # Add a textbox with hyperlink text
            textbox = slide.shapes.add_textbox(x, y, Inches(5), Inches(0.5))
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = url
            run.font.color.rgb = PPTXColor(0, 0, 255)  # Blue color for hyperlink

        prs.save(file_path)
        return file_path

    @pytest.fixture(autouse=True)
    def setup(self):
        # Folder path for test files
        self.folder_path = "tests/test_files"
        self.pdf_file_path = self.create_url_only_pdf_file(self.folder_path, "url_only_file.pdf")
        self.docx_file_path = self.create_url_only_docx_file(self.folder_path, "url_only_file.docx")
        self.pptx_file_path = self.create_url_only_pptx_file(self.folder_path, "url_only_file.pptx")

    def test_text_only_pdf_file_loading(self):
        # Initialize the PDFLoader with the path to the test PDF file
        pdf_loader = PDFLoader(self.pdf_file_path)
        pdf_loader.validate_file()

        # Use PDFLoader's load_file method to open and read the PDF
        pdf_document = pdf_loader.load_file()
        assert pdf_document is not None, "Failed to load the url only PDF document."

        # Extract text from the PDF document
        extracted_text = ""
        for page in pdf_document:
            extracted_text += page.get_text()

        # Display the loaded PDF text
        print("Loaded PDF Text Content:\n")
        print(extracted_text)  # Display the entire extracted text on terminal

    def test_text_only_docx_file_loading(self):
        #Initialize the DOCXLoader with the path to the test DOCX file
        docx_loader = DOCXLoader(self.docx_file_path)
        docx_loader.validate_file()
        #Use DOCXLoader's load_file method to load the docx
        docx_document = docx_loader.load_file()
        
        assert docx_document is not None, "Failed to load the url only DOCX document"

        # Extract text from the Document instance
        extracted_text = []
        for paragraph in docx_document.paragraphs:
            extracted_text.append(paragraph.text)

        # Combine the paragraphs into a single string
        full_text = "\n".join(extracted_text)

        # Display the loaded DOCX text
        print("Loaded DOCX Text Content:\n")
        print(full_text)  # Display the entire extracted text

    def test_text_only_pptx_file_loading(self):
        #Initialize the PPTXLoader with the path to the test PPTX file
        pptx_loader = PPTLoader(self.pptx_file_path)
        pptx_loader.validate_file()
        #Use PPTLoader's load_file method to load the pptx
        pptx_document = pptx_loader.load_file()

        assert pptx_document is not None, "Failed to load the url only PPTX document."

        # Extract text from the presentation
        extracted_text = []
        for slide in pptx_document.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Ensure the shape has text
                    extracted_text.append(shape.text)

        # Combine the text from all slides into a single string
        full_text = "\n".join(extracted_text)

        # Display the loaded PPTX text
        print("Loaded PPTX Text Content:\n")
        print(full_text)  # Display the entire extracted text